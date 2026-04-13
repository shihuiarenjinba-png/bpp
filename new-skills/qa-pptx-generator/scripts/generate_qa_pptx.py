import sys
import json
from pptx import Presentation
import os
import copy

def fill_slide_content(slide, slide_type, data):
    """Fills content into a specific slide based on its type (question/answer)."""
    q_data = data['question']
    a_data = data['answer']
    
    def set_shape_text(slide, shape_id, text):
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                shape.text = text
                return True
        return False

    if slide_type == 'question':
        set_shape_text(slide, 6, f"問題：{q_data.get('category', '一般')}")
        set_shape_text(slide, 7, f"Q. {q_data['text']}")
        set_shape_text(slide, 8, q_data['choices'][0])
        set_shape_text(slide, 10, q_data['choices'][1])
        set_shape_text(slide, 12, q_data['choices'][2])
        set_shape_text(slide, 14, q_data['choices'][3])
    elif slide_type == 'answer':
        set_shape_text(slide, 6, f"回答：{q_data.get('category', '一般')}")
        set_shape_text(slide, 2, f"A. {a_data['correct_choice']}")
        set_shape_text(slide, 4, a_data['explanation'])

def duplicate_slide(prs, index):
    """
    Duplicates a slide within the same presentation, including background and all shapes.
    """
    source_slide = prs.slides[index]
    # Use the same layout to ensure the slide structure is consistent
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    
    # 1. Copy Background (Fixing the grey background issue)
    # We use XML manipulation to copy the background element from the source slide
    try:
        # Define namespaces for finding elements
        nsmap = source_slide._element.nsmap
        # Find the background element in the source slide
        source_bg = source_slide._element.find('.//p:bg', namespaces=nsmap)
        if source_bg is not None:
            # Deep copy the background XML
            new_bg = copy.deepcopy(source_bg)
            # Remove any existing background in the new slide
            existing_bg = new_slide._element.find('.//p:bg', namespaces=nsmap)
            if existing_bg is not None:
                new_slide._element.remove(existing_bg)
            # Insert the copied background at the beginning of the slide XML
            new_slide._element.insert(0, new_bg)
    except Exception as e:
        print(f"Warning: Failed to copy background XML: {e}")

    # 2. Remove default shapes from the newly created slide to avoid duplicates
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape._element)
    
    # 3. Copy all shapes from the template slide's XML element to the new slide
    for shape in source_slide.shapes:
        new_element = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_element)
    
    return new_slide

def generate_pptx(data_list, template_path, output_path):
    # Load the template
    prs = Presentation(template_path)
    
    # Process each quiz
    for i, qa_item in enumerate(data_list):
        if i == 0:
            # First quiz: Use the existing slides
            fill_slide_content(prs.slides[0], 'question', qa_item)
            fill_slide_content(prs.slides[1], 'answer', qa_item)
        else:
            # Subsequent quizes: Duplicate the first two slides (Question & Answer)
            new_q_slide = duplicate_slide(prs, 0)
            new_a_slide = duplicate_slide(prs, 1)
            
            fill_slide_content(new_q_slide, 'question', qa_item)
            fill_slide_content(new_a_slide, 'answer', qa_item)

    prs.save(output_path)
    print(f"PPTX with {len(data_list)} quizes saved to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python generate_qa_pptx.py <json_data_path> <output_pptx_path>")
        sys.exit(1)
    
    json_path = sys.argv[1]
    out_path = sys.argv[2]
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_pptx = os.path.join(script_dir, "../templates/template.pptx")
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    if isinstance(data, dict):
        data_list = [data]
    else:
        data_list = data
    
    generate_pptx(data_list, template_pptx, out_path)
