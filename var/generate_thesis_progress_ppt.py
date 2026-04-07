from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from thesis_research_state import (
    BACKGROUND,
    CHAPTER_STATUS,
    CONTRIBUTIONS,
    CURRENT_ISSUES,
    CURRENT_PROGRESS,
    DATA_PLAN,
    HYPOTHESES,
    PPT_OUTPUT_PATH,
    PURPOSES,
    ROADMAP,
    THESIS_TITLE,
    THREE_LAYER_STRUCTURE,
)


OUT_PATH = Path(PPT_OUTPUT_PATH)

SLIDE_W = 1280
SLIDE_H = 720

BLACK = RGBColor(26, 26, 26)
DARK = RGBColor(68, 68, 68)
MID = RGBColor(118, 118, 118)
LIGHT = RGBColor(214, 214, 214)
PANEL = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)

TITLE_SIZE = 22.5
SUB_SIZE = 19.5
BODY_SIZE = 12.5
SMALL_SIZE = 10.0


def px(value: int) -> Emu:
    return Emu(value * 9525)


def add_text(slide, x, y, w, h, text, size, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "IPAexGothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, size=BODY_SIZE, color=BLACK, gap=8):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.space_after = Pt(gap)
        for run in p.runs:
            run.font.name = "IPAexGothic"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_panel(slide, x, y, w, h, fill=PANEL, line=LIGHT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px(x), px(y), px(w), px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def draw_line(slide, x1, y1, x2, y2, color=DARK, width=1.0):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px(x1), px(y1), px(max(1, x2 - x1)), px(max(1, y2 - y1)))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_frame(slide, title, subtitle, page_no):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_text(slide, 20, 16, 1220, 34, title, TITLE_SIZE, BLACK, True)
    add_text(slide, 20, 80, 1220, 28, subtitle, SUB_SIZE, DARK, False)
    draw_line(slide, 13, 67, 1260, 69, DARK)
    draw_line(slide, 13, 688, 1260, 690, LIGHT)
    add_text(slide, 20, 694, 930, 18, THESIS_TITLE, SMALL_SIZE, MID)
    add_text(slide, 1180, 692, 60, 18, str(page_no), SMALL_SIZE, MID, align=PP_ALIGN.RIGHT)


def build_deck():
    prs = Presentation()
    prs.slide_width = px(SLIDE_W)
    prs.slide_height = px(SLIDE_H)
    layout = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(layout)
    add_frame(slide, "研究の核心", "最終計画書ベースで、研究の狙いと現在地だけを整理", 1)
    add_panel(slide, 20, 120, 600, 230)
    add_text(slide, 36, 138, 560, 24, "問題意識", 15, BLACK, True)
    add_bullets(slide, 36, 170, 550, 155, BACKGROUND, 12.2)
    add_panel(slide, 650, 120, 610, 230)
    add_text(slide, 666, 138, 560, 24, "研究の狙い", 15, BLACK, True)
    add_bullets(slide, 666, 170, 560, 160, [
        "機械学習の高い個別株予測力を、そのまま使うのではなく、なぜ当たるのかを因果性で説明する。",
        "Fama-Frenchファクターとマクロ経済変数の先行・遅行構造を明らかにする。",
        "主データは最終計画書で、他資料は補助的な参考情報として扱う。",
    ], 12.2)
    add_panel(slide, 20, 382, 1240, 268, fill=WHITE)
    add_text(slide, 36, 400, 1180, 24, "現在の位置づけ", 15, BLACK, True)
    add_bullets(slide, 36, 432, 1180, 190, [
        "研究の主軸は『予測精度の競争』ではなく、『予測力の源泉を経済的に解釈すること』に置く。",
        "当初の4層ブロック案は棄却され、現在は 3層構造で再構築する方針が最終計画書で明示されている。",
        "最大の未了項目は、2006年以降 239か月より長い正規データを取得して、再走すること。",
    ], 12.4)

    # 2
    slide = prs.slides.add_slide(layout)
    add_frame(slide, "研究目的と仮説", "目的・仮説を1枚で把握できるように分離", 2)
    add_panel(slide, 20, 120, 560, 530)
    add_text(slide, 36, 138, 520, 24, "研究目的", 15, BLACK, True)
    add_bullets(slide, 36, 170, 510, 460, PURPOSES, 12.2)
    xs = [620, 835, 1050]
    widths = [195, 195, 210]
    for idx, (hyp, x, w) in enumerate(zip(HYPOTHESES, xs, widths), start=1):
        add_panel(slide, x, 120, w, 530, fill=PANEL)
        add_text(slide, x + 12, 138, w - 24, 22, hyp["id"], 14, BLACK, True, PP_ALIGN.CENTER)
        add_text(slide, x + 12, 172, w - 24, 58, hyp["title"], 14, BLACK, True, PP_ALIGN.CENTER)
        add_text(slide, x + 12, 248, w - 24, 360, hyp["detail"], 11.6, DARK)

    # 3
    slide = prs.slides.add_slide(layout)
    add_frame(slide, "先行研究と研究ギャップ", "PowerPoint では差別化の論点だけに絞る", 3)
    add_panel(slide, 20, 120, 400, 530)
    add_text(slide, 36, 138, 360, 24, "ML側の先行研究", 15, BLACK, True)
    add_bullets(slide, 36, 170, 350, 450, [
        "Gu, Kelly, Xiu (2020) では、木系モデルとニューラルネットが高い個別株予測力を示した。",
        "ただし、その予測力の背後にある経済的な因果メカニズムは直接説明していない。",
        "本研究は、このブラックボックス問題を計量経済学的に掘り下げる。",
    ], 12.2)
    add_panel(slide, 440, 120, 400, 530)
    add_text(slide, 456, 138, 360, 24, "FF・マクロ側の先行研究", 15, BLACK, True)
    add_bullets(slide, 456, 170, 350, 450, [
        "Fama-Frenchファクターとマクロ変数の関連は多く研究されてきた。",
        "ただし、双方向の先行・遅行構造を主題に据えて、個別株予測の源泉まで接続する研究は限られる。",
        "CPI→RMW のような時系列的な因果性を、重点発見として扱う余地がある。",
    ], 12.2)
    add_panel(slide, 860, 120, 400, 530)
    add_text(slide, 876, 138, 360, 24, "本研究の差別化", 15, BLACK, True)
    add_bullets(slide, 876, 170, 350, 450, [
        "MLの高い予測力を『使う』のではなく『説明する』方向へ寄せる。",
        "FFファクターとマクロ経済の双方向因果性を、VARとグレンジャー因果性で検証する。",
        "構造の可視化、CPI→RMW、t分布残差を組み合わせて学術的・実務的示唆へつなげる。",
    ], 12.2)

    # 4
    slide = prs.slides.add_slide(layout)
    add_frame(slide, "3層構造と主要変数", "最終計画書で明示された層構造と変数の置き方", 4)
    x_positions = [30, 448, 866]
    widths = [374, 374, 374]
    for item, x, w in zip(THREE_LAYER_STRUCTURE, x_positions, widths):
        add_panel(slide, x, 150, w, 290)
        add_text(slide, x + 14, 168, w - 28, 24, item["layer"], 15, BLACK, True, PP_ALIGN.CENTER)
        add_text(slide, x + 14, 208, w - 28, 44, item["variables"], 13.2, DARK, True, PP_ALIGN.CENTER)
        add_text(slide, x + 16, 268, w - 32, 140, item["interpretation"], 11.8, MID, False, PP_ALIGN.CENTER)
    add_panel(slide, 20, 472, 1240, 178, fill=WHITE)
    add_text(slide, 36, 490, 1180, 24, "データ方針", 15, BLACK, True)
    add_bullets(slide, 36, 522, 1180, 110, DATA_PLAN, 12.2)

    # 5
    slide = prs.slides.add_slide(layout)
    add_frame(slide, "分析フロー", "1枚1話題に合わせて、実証の順番だけを固定化", 5)
    stages = [
        ("Step 1", "正規データの取得と再走"),
        ("Step 2", "前処理: 差分化・標準化・スケール調整"),
        ("Step 3", "3層ブロック外生性VARの再構築"),
        ("Step 4", "ラグ選択とグレンジャー因果性検定"),
        ("Step 5", "多重検定補正と重点関係の抽出"),
        ("Step 6", "IRF と FEVD による動的分析"),
        ("Step 7", "OOS R² と方向的中率の評価"),
        ("Step 8", "t分布残差によるテールリスク確認"),
    ]
    y = 134
    for step, text in stages:
        add_panel(slide, 32, y, 1216, 54, fill=PANEL)
        add_text(slide, 50, y + 14, 120, 20, step, 13, BLACK, True)
        add_text(slide, 178, y + 14, 1000, 20, text, 13, DARK)
        y += 62

    # 6
    slide = prs.slides.add_slide(layout)
    add_frame(slide, "データ・検定・評価の着眼点", "Excel 側で深掘りする内容を、スライドでは要点だけ示す", 6)
    add_panel(slide, 20, 120, 400, 530)
    add_text(slide, 36, 138, 360, 24, "データ処理", 15, BLACK, True)
    add_bullets(slide, 36, 170, 350, 430, [
        "非定常なスプレッド変数は一次差分または対数差分を検討する。",
        "VIX は対数差分化または標準化でスケールを調整する。",
        "長期月次データを使い、リーマンショック以前を含めて頑健性を高める。",
    ], 12.2)
    add_panel(slide, 440, 120, 400, 530)
    add_text(slide, 456, 138, 360, 24, "重点検定", 15, BLACK, True)
    add_bullets(slide, 456, 170, 350, 430, [
        "CPI→RMW は 1, 3, 6, 12か月など複数ラグで重点確認する。",
        "逆方向の因果性も検証し、双方向構造の有無を明示する。",
        "Bonferroni などの補正後も残る関係を強い発見として扱う。",
    ], 12.2)
    add_panel(slide, 860, 120, 400, 530)
    add_text(slide, 876, 138, 360, 24, "評価指標", 15, BLACK, True)
    add_bullets(slide, 876, 170, 350, 430, [
        "IRF でショックの伝播経路を確認する。",
        "FEVD で変数間の相対的重要性を整理する。",
        "OOS R² と方向的中率は補助的な実務評価として用いる。",
        "t分布残差でファットテールを反映したリスク評価を行う。",
    ], 12.2)

    # 7
    slide = prs.slides.add_slide(layout)
    add_frame(slide, "現在の進捗とボトルネック", "何が固まっていて、何がまだ未了かを分けて表示", 7)
    add_panel(slide, 20, 120, 620, 260)
    add_text(slide, 36, 138, 580, 24, "固まっていること", 15, BLACK, True)
    add_bullets(slide, 36, 170, 560, 180, CURRENT_PROGRESS, 12.1)
    add_panel(slide, 660, 120, 600, 260)
    add_text(slide, 676, 138, 560, 24, "ボトルネック", 15, BLACK, True)
    add_bullets(slide, 676, 170, 540, 180, CURRENT_ISSUES, 12.1)
    add_panel(slide, 20, 400, 1240, 250, fill=WHITE)
    add_text(slide, 36, 418, 1180, 24, "章進捗", 15, BLACK, True)
    y = 456
    for chapter, title, status, note in CHAPTER_STATUS:
        add_text(slide, 40, y, 90, 18, chapter, 11.5, BLACK, True)
        add_text(slide, 136, y, 190, 18, title, 11.5, DARK, True)
        add_text(slide, 332, y, 90, 18, status, 11.5, MID, True)
        add_text(slide, 430, y, 790, 18, note, 11.0, DARK)
        y += 30

    # 8
    slide = prs.slides.add_slide(layout)
    add_frame(slide, "ロードマップ", "直近で動く順番と成果物の使い分けを明示", 8)
    add_panel(slide, 20, 120, 860, 530)
    add_text(slide, 36, 138, 820, 24, "直近の実行順", 15, BLACK, True)
    add_bullets(slide, 36, 170, 790, 450, ROADMAP, 12.1)
    add_panel(slide, 900, 120, 360, 530)
    add_text(slide, 916, 138, 320, 24, "成果物の役割分担", 15, BLACK, True)
    add_bullets(slide, 916, 170, 300, 220, [
        "PowerPoint: 研究の全体像を短時間で説明するための資料。",
        "Excel: 文献・変数・データソース・検定設計を蓄積する研究ベース。",
        "コードファイル: 現在の前提、仮説、ロードマップをコードとして固定する管理用ファイル。",
    ], 11.8)
    add_text(slide, 916, 430, 300, 110, "最優先は、長期の正規データ取得と再走です。ここが終わると、因果性検定と図表作成が一気に進めやすくなります。", 12.0, DARK)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)


if __name__ == "__main__":
    build_deck()
    print(OUT_PATH)
